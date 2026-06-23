"""PPTX extractor: text per slide (tables, notes) + image understanding."""

import hashlib
import io
from collections.abc import Iterable
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app.core.errors import ExtractionError
from app.ingestion.detector import register_extractor
from app.ingestion.models import Element, ElementType, ExtractedDocument, Location
from app.ingestion.vision import images_to_elements

_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


def _table_text(table: Any) -> str:
    rows = [" | ".join(cell.text.strip() for cell in row.cells) for row in table.rows]
    return "\n".join(r for r in rows if r.strip(" |"))


def _collect(shapes: Iterable[Any], texts: list[str], images: list[tuple[bytes, str]]) -> None:
    """Walk shapes (recursing into groups); gather text and picture blobs."""
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            _collect(shape.shapes, texts, images)
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            images.append((shape.image.blob, shape.image.content_type))
        elif getattr(shape, "has_table", False):
            texts.append(_table_text(shape.table))
        elif shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                texts.append(text)


def extract_pptx(filename: str, data: bytes) -> ExtractedDocument:
    try:
        prs = Presentation(io.BytesIO(data))
    except Exception as exc:
        raise ExtractionError(f"Failed to open PPTX: {exc}") from exc

    elements: list[Element] = []
    images: list[tuple[bytes, str, Location]] = []

    for slide_no, slide in enumerate(prs.slides, start=1):
        texts: list[str] = []
        slide_images: list[tuple[bytes, str]] = []
        _collect(slide.shapes, texts, slide_images)

        slide_text = "\n".join(texts).strip()
        if slide_text:
            elements.append(
                Element(
                    text=slide_text,
                    element_type=ElementType.SLIDE_TEXT,
                    location=Location(slide=slide_no),
                )
            )
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                elements.append(
                    Element(
                        text=notes,
                        element_type=ElementType.NOTES,
                        location=Location(slide=slide_no),
                    )
                )
        images.extend((blob, mime, Location(slide=slide_no)) for blob, mime in slide_images)

    elements.extend(images_to_elements(images))

    if not elements:
        raise ExtractionError("PPTX has no extractable text")

    return ExtractedDocument(
        filename=filename,
        mime_type=_MIME,
        format="pptx",
        file_hash=hashlib.sha256(data).hexdigest(),
        size_bytes=len(data),
        elements=elements,
        num_units=len(prs.slides),
    )


register_extractor("pptx", extract_pptx)
