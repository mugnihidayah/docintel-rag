import base64
import io
from types import SimpleNamespace

import pytest
from pptx import Presentation
from pptx.util import Inches

import app.ingestion.vision as vision_mod
from app.ingestion.detector import get_extractor
from app.ingestion.models import ElementType
from app.ingestion.pptx import extract_pptx

_PNG_1X1 = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mNk+M9QDwADhgGAWjR9awAAAABJRU5ErkJggg=="
)


def _make_pptx() -> bytes:
    prs = Presentation()
    s1 = prs.slides.add_slide(prs.slide_layouts[1])
    s1.shapes.title.text = "Welcome" # type: ignore
    s1.placeholders[1].text = "First slide body text" # type: ignore
    s1.notes_slide.notes_text_frame.text = "Remember to smile" # type: ignore
    s2 = prs.slides.add_slide(prs.slide_layouts[5])
    s2.shapes.title.text = "Second Slide" # type: ignore
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _make_pptx_with_picture() -> bytes:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Charts" # type: ignore
    slide.shapes.add_picture(io.BytesIO(_PNG_1X1), Inches(1), Inches(1), Inches(2), Inches(2))
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def test_extract_pptx_slides_and_notes() -> None:
    result = extract_pptx("deck.pptx", _make_pptx())
    assert result.format == "pptx"
    assert result.num_units == 2
    slide_texts = [e for e in result.elements if e.element_type == ElementType.SLIDE_TEXT]
    assert any("Welcome" in e.text and e.location.slide == 1 for e in slide_texts)
    assert any("Second Slide" in e.text and e.location.slide == 2 for e in slide_texts)
    notes = [e for e in result.elements if e.element_type == ElementType.NOTES]
    assert any("Remember to smile" in e.text and e.location.slide == 1 for e in notes)


def test_extract_pptx_image_via_vision(monkeypatch: pytest.MonkeyPatch) -> None:
    monkeypatch.setattr(vision_mod, "describe_image", lambda blob, mime: "chart showing growth")
    monkeypatch.setattr(
        vision_mod,
        "get_settings",
        lambda: SimpleNamespace(vision_min_image_kb=0, vision_max_images=10),
    )
    result = extract_pptx("deck.pptx", _make_pptx_with_picture())
    image_els = [e for e in result.elements if e.element_type == ElementType.IMAGE]
    assert any("chart showing growth" in e.text for e in image_els)
    assert image_els[0].location.slide == 1


def test_pptx_extractor_is_registered() -> None:
    assert get_extractor("pptx") is extract_pptx
